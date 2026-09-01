from __future__ import annotations

import base64
from html import escape
from io import BytesIO
from pathlib import Path
from typing import Iterable
from urllib.parse import urlparse

import httpx
from PIL import Image, ImageDraw, ImageFont, ImageOps
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfbase.ttfonts import TTFError, TTFont
from reportlab.pdfgen import canvas

from ..core.config import get_settings
from ..domain.models import Trip


class ReportAgent:
    """Build every export from one frozen, verified Trip snapshot."""

    def render(self, trip: Trip, markdown: str, kind: str) -> bytes:
        lines = export_lines(trip, markdown)
        renderer = {
            "pdf": render_pdf,
            "pptx": render_pptx,
            "png": render_long_image,
        }.get(kind)
        if kind == "html":
            return build_report_html(trip, markdown).encode("utf-8")
        if renderer is None:
            raise ValueError(f"unsupported report format: {kind}")
        return renderer(lines, trip)


def build_report_html(trip: Trip, markdown: str = "") -> str:
    """Create the single visual report template shared by export consumers.

    PDF/PPTX/PNG keep native rendering for offline reliability, while this
    template is also available as a standalone HTML export and defines the
    same hierarchy: cover route, curated cards, then day-by-day details.
    """
    route_uri = _png_data_uri(_render_route_map(trip, 1440, 560).getvalue())
    seen_image_urls: set[str] = set()
    all_activities = [
        activity
        for day in trip.days
        for activity in day.activities
        if activity.type in {"attraction", "meal", "hotel"}
    ]

    def card(activity) -> str:
        image = _activity_image(activity, seen_image_urls)
        image_markup = (
            f'<img src="{_image_data_uri(_image_bytes(image))}" alt="{escape(activity.place.name)}">'
            if image else "<div class=\"report-card-placeholder\">RoadMan 精选</div>"
        )
        description = escape(activity.description or activity.user_note or "详情见来源链接")
        detail = escape(activity.detail_url or "")
        link = f'<a href="{detail}" target="_blank" rel="noreferrer">查看详情与来源</a>' if detail else ""
        reservation = {
            "required": "需预约",
            "recommended": "建议预约",
            "not_required": "无需预约",
            "unknown": "预约待核查",
        }.get(activity.reservation_status, "预约待核查")
        checks = (
            f'<div class="report-checks"><span>{reservation}</span>'
            + "".join(f"<span>{escape(tag)}</span>" for tag in activity.risk_tags[:3])
            + "</div>"
        )
        ticket = activity.ticket_or_price
        if ticket:
            ticket_text = f"¥{ticket.minimum:g}"
            if ticket.maximum != ticket.minimum:
                ticket_text += f"–¥{ticket.maximum:g}"
        elif activity.ticket_status == "free":
            ticket_text = "免费"
        else:
            ticket_text = "暂未返回"
        parking = activity.parking_note or "暂未返回"
        if activity.parking_or_price:
            parking = f"约¥{activity.parking_or_price.minimum:g}"
            if activity.parking_or_price.maximum != activity.parking_or_price.minimum:
                parking += f"–¥{activity.parking_or_price.maximum:g}"
            if activity.parking_note:
                parking += f"（{activity.parking_note}）"
        facts = [
            f"营业：{activity.opening_hours.text if activity.opening_hours else '暂未返回'}",
            f"门票：{ticket_text}",
            f"停车：{parking}",
            f"资料：{'完整' if activity.information_status == 'complete' else '部分' if activity.information_status == 'partial' else '待核验'}（{activity.information_sources_count} 个来源）",
        ]
        facts_markup = '<div class="report-facts">' + "".join(
            f"<span>{escape(item)}</span>" for item in facts
        ) + "</div>"
        return (
            f'<article class="report-card"><div class="report-card-media">{image_markup}'
            f'<span class="report-card-type">{escape(_activity_label(activity.type))}</span></div>'
            f'<h3>{escape(activity.place.name)}</h3>'
            f'<time>{activity.planned_start:%m月%d日 %H:%M}–{activity.planned_end:%H:%M}</time>'
            f'<p>{description[:180]}</p>{facts_markup}{checks}{link}</article>'
        )

    day_sections = []
    for day in trip.days:
        stage_items = []
        for stage in day.stages:
            service = ""
            if stage.mode in {"train", "flight", "ferry"}:
                terminals = " → ".join(
                    item for item in (stage.departure_terminal, stage.arrival_terminal) if item
                )
                service = f" · {escape(stage.service_number or '班次号暂未返回')}"
                if terminals:
                    service += f" · {escape(terminals)}"
            transit = ""
            if stage.transit_legs:
                transit = " · " + "；".join(
                    f"{escape(leg.line_name or leg.line_type or '公共交通')} "
                    f"{escape(leg.departure_stop or '上车')}→{escape(leg.arrival_stop or '下车')}"
                    for leg in stage.transit_legs
                )
            stage_items.append(
                f'<li><b>{stage.planned_start:%H:%M}–{stage.planned_end:%H:%M}</b> '
                f'{escape(stage.origin.name)} → {escape(stage.destination.name)} · '
                f'{stage.distance_km:g} km · {_mode_label(stage.mode)}{service}{transit}</li>'
            )
        stages = "".join(stage_items)
        cards = "".join(card(activity) for activity in day.activities if activity.type in {"attraction", "meal", "hotel"})
        day_route = _png_data_uri(_render_route_map(trip, 1100, 300, days=[day]).getvalue())
        day_sections.append(
            f'<section class="report-day"><div class="report-day-heading"><span>第 {day.day_index} 天</span>'
            f'<h2>{escape(day.title)}</h2><small>{day.date} · {day.total_distance_km:g} km · '
            f'{_format_duration(day.total_drive_minutes)}驾驶</small></div>'
            f'<img class="day-route" src="{day_route}" alt="第 {day.day_index} 天路线图">'
            f'<ol class="report-stages">{stages}</ol><div class="report-cards">{cards}</div></section>'
        )
    cards = "".join(card(activity) for activity in all_activities[:4])
    return f"""<!doctype html>
<html lang="zh-CN"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>{escape(trip.title)} · RoadMan 行程报告</title>
<style>
:root {{ color-scheme: light; font-family: Inter,"Microsoft YaHei",sans-serif; color:#10213e; background:#edf4fc; }}
* {{ box-sizing:border-box; }} body {{ margin:0; padding:36px; }}
.report {{ max-width:1180px; margin:auto; }}
.cover {{ padding:42px; border-radius:30px; color:white; background:linear-gradient(135deg,#10213e,#1e5faa); box-shadow:0 20px 50px #183d7040; }}
.eyebrow {{ letter-spacing:.16em; text-transform:uppercase; color:#8ee4f4; font-size:13px; font-weight:800; }}
h1 {{ margin:16px 0 8px; font-size:42px; }} .subtitle {{ color:#d7e8fb; font-size:17px; }}
.route {{ width:100%; margin-top:28px; border-radius:20px; background:#f3f8ff; }}
.section {{ margin-top:28px; padding:30px; border-radius:26px; background:#fff; box-shadow:0 12px 34px #1b47701c; }}
.section h2 {{ margin:0 0 18px; font-size:24px; }} .report-cards {{ display:grid; grid-template-columns:repeat(4,minmax(0,1fr)); gap:18px; }}
.report-card {{ overflow:hidden; border:1px solid #d9e6f4; border-radius:18px; background:#fff; }}
.report-card-media {{ position:relative; height:170px; background:#edf4ff; }} .report-card-media img {{ width:100%; height:100%; object-fit:cover; }}
.report-card-placeholder {{ display:grid; place-items:center; height:100%; color:#2377e8; font-weight:800; font-size:20px; }}
.report-card-type {{ position:absolute; left:12px; top:12px; padding:5px 9px; border-radius:999px; color:#fff; background:#2377e8dd; font-size:12px; font-weight:800; }}
 .report-card h3 {{ margin:14px 14px 5px; font-size:17px; }} .report-card time,.report-card p,.report-card a {{ display:block; margin:0 14px 9px; color:#657b98; font-size:13px; line-height:1.5; }} .report-card a {{ color:#176fe1; text-decoration:none; font-weight:700; }}
.report-facts {{ display:grid; gap:3px; margin:0 14px 9px; color:#526b89; font-size:11px; line-height:1.35; }}
.report-checks {{ display:flex; flex-wrap:wrap; gap:5px; margin:0 14px 9px; }} .report-checks span {{ padding:3px 7px; border-radius:999px; color:#7d5d21; background:#fff4d8; font-size:11px; }}
.report-day {{ margin-top:28px; padding:26px; border:1px solid #d9e6f4; border-radius:24px; background:#f8fbff; }}
.day-route {{ width:100%; margin:8px 0 18px; border-radius:16px; background:#eef5fc; }}
.report-day-heading {{ display:flex; align-items:baseline; gap:12px; flex-wrap:wrap; }} .report-day-heading span {{ color:#176fe1; font-weight:800; }} .report-day-heading h2 {{ margin:0; font-size:25px; }} .report-day-heading small {{ color:#7186a0; }}
.report-stages {{ margin:18px 0; padding-left:25px; color:#385a80; }} .report-stages li {{ margin:8px 0; }} .report-stages b {{ color:#10213e; margin-right:10px; }}
@media(max-width:800px) {{ body {{ padding:16px; }} .report-cards {{ grid-template-columns:repeat(2,minmax(0,1fr)); }} h1 {{ font-size:30px; }} }}
</style></head><body><main class="report">
<header class="cover"><div class="eyebrow">ROADMAN · 智能行程报告</div><h1>{escape(trip.title)}</h1>
<div class="subtitle">{escape(str(trip.start_date or "待定"))} – {escape(str(trip.end_date or "待定"))} · {escape(trip.origin.name if trip.origin else "待定")} → {escape(trip.destination.name if trip.destination else "待定")}</div>
<img class="route" src="{route_uri}" alt="行程路线图"></header>
<section class="section"><h2>景点、餐饮与住宿</h2><div class="report-cards">{cards}</div></section>
{"".join(day_sections)}
</main></body></html>"""


def _png_data_uri(data: bytes) -> str:
    return "data:image/png;base64," + base64.b64encode(data).decode("ascii")


def _image_data_uri(data: bytes) -> str:
    return "data:image/jpeg;base64," + base64.b64encode(data).decode("ascii")


def _image_bytes(image: Image.Image | None) -> bytes:
    output = BytesIO()
    if image is None:
        return b""
    image = image.copy()
    image.thumbnail((960, 640), Image.Resampling.LANCZOS)
    image.save(output, format="JPEG", quality=82, optimize=True)
    return output.getvalue()


def _activity_label(value: str) -> str:
    return {"attraction": "景点", "hotel": "住宿", "meal": "餐饮"}.get(value, "安排")


def _mode_label(value: str) -> str:
    return {"driving": "驾车", "transit": "公共交通", "walking": "步行", "riding": "骑行"}.get(value, value)


def _format_duration(minutes: int) -> str:
    hours, rest = divmod(max(0, int(minutes)), 60)
    if hours and rest:
        return f"{hours}小时{rest}分钟"
    return f"{hours}小时" if hours else f"{rest}分钟"


def export_lines(trip: Trip, markdown: str) -> list[str]:
    """Build one deterministic snapshot used by every export format."""
    lines = [
        trip.title,
        f"日期：{trip.start_date or '待定'} 至 {trip.end_date or '待定'}",
        f"路线：{trip.origin.name if trip.origin else '待定'} → {trip.destination.name if trip.destination else '待定'}",
        "",
    ]
    if markdown:
        lines.extend(markdown.splitlines())
    else:
        for day in trip.days:
            lines.append(f"第 {day.day_index} 天 · {day.date} · {day.title}")
            for stage in day.stages:
                lines.append(
                    f"{stage.origin.name} → {stage.destination.name} · "
                    f"{stage.planned_start:%H:%M}-{stage.planned_end:%H:%M} · "
                    f"{stage.mode} · {stage.distance_km:g} km"
                )
            for activity in day.activities:
                lines.append(
                    f"{activity.planned_start:%H:%M}-{activity.planned_end:%H:%M} "
                    f"{activity.type} · {activity.place.name}"
                )
            lines.append("")
    return lines


def _register_report_font() -> str:
    """Register a real CJK TrueType font for PDF text.

    ReportLab's CID fallback can render Chinese inconsistently in browser PDF
    viewers (missing glyphs or blank headings). Prefer the fonts installed in
    the container/Windows host and keep the CID font as a last-resort fallback.
    """
    candidates = (
        "/usr/share/fonts/truetype/arphic/ukai.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "C:/Windows/Fonts/msyh.ttc",
        "C:/Windows/Fonts/msyh.ttf",
    )
    for index, candidate in enumerate(candidates):
        if not Path(candidate).is_file():
            continue
        try:
            name = f"RoadManCJK{index}"
            pdfmetrics.registerFont(TTFont(name, candidate, subfontIndex=0))
            return name
        except (OSError, TTFError, ValueError, IndexError):
            continue
    pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
    return "STSong-Light"


def render_pdf(lines: Iterable[str], trip: Trip | None = None) -> bytes:
    output = BytesIO()
    font_name = _register_report_font()
    document = canvas.Canvas(output, pagesize=(595, 842))
    document.setTitle("RoadMan 行程安排")
    y = 800
    if trip and trip.days:
        # Cover: a visual route summary with a restrained report hierarchy.
        document.setFillColorRGB(0.06, 0.13, 0.25)
        document.rect(0, 0, 595, 842, fill=1, stroke=0)
        document.setFillColorRGB(0.25, 0.60, 0.96)
        document.setFont(font_name, 11)
        document.drawString(42, 778, "ROADMAN  ·  智能行程报告")
        document.setFillColorRGB(1, 1, 1)
        document.setFont(font_name, 27)
        document.drawString(42, 730, trip.title[:25])
        document.setFont(font_name, 12)
        document.setFillColorRGB(0.76, 0.85, 0.95)
        document.drawString(
            42,
            700,
            f"{trip.start_date or '待定'} — {trip.end_date or '待定'}   ·   {trip.origin.name if trip.origin else '待定'} → {trip.destination.name if trip.destination else '待定'}",
        )
        route_image = _render_route_map(trip, 1200, 520)
        document.drawImage(ImageReader(route_image), 42, 430, width=511, height=221, preserveAspectRatio=True)
        activity_image = _render_activity_board(trip, 1200, 300)
        document.drawImage(ImageReader(activity_image), 42, 75, width=511, height=128, preserveAspectRatio=True)
        document.showPage()
        for day_index, day in enumerate(trip.days):
            document.setFillColorRGB(0.96, 0.98, 1)
            document.rect(0, 0, 595, 842, fill=1, stroke=0)
            document.setFillColorRGB(0.06, 0.13, 0.25)
            document.setFont(font_name, 21)
            document.drawString(42, 790, f"第 {day.day_index} 天  ·  {day.date}  {day.title}")
            document.setFillColorRGB(0.28, 0.39, 0.55)
            document.setFont(font_name, 11)
            document.drawString(42, 766, f"总里程 {day.total_distance_km:g} km   ·   驾驶 {day.total_drive_minutes // 60}h{day.total_drive_minutes % 60}m")
            day_route = _render_route_map(trip, 1200, 300, days=[day])
            document.drawImage(ImageReader(day_route), 42, 505, width=511, height=128, preserveAspectRatio=True)
            day_image = _render_activity_board(trip, 1200, 420, activities=day.activities)
            document.drawImage(ImageReader(day_image), 42, 315, width=511, height=179, preserveAspectRatio=True)
            y = 275
            document.setFillColorRGB(0.10, 0.20, 0.34)
            for stage in day.stages:
                line = f"{stage.planned_start:%H:%M}–{stage.planned_end:%H:%M}   {stage.origin.name} → {stage.destination.name}   {stage.distance_km:g} km"
                for wrapped in _wrap(line, 54):
                    document.drawString(42, y, wrapped)
                    y -= 18
                if y < 70:
                    break
            if day_index < len(trip.days) - 1:
                document.showPage()
        document.save()
        return output.getvalue()
    for index, line in enumerate(lines):
        document.setFont(font_name, 18 if index == 0 else 10)
        for wrapped in _wrap(line, 28 if index == 0 else 48):
            if y < 45:
                document.showPage()
                document.setFont(font_name, 10)
                y = 800
            document.drawString(42, y, wrapped)
            y -= 18 if index == 0 else 15
        if index == 0:
            y -= 8
    document.save()
    return output.getvalue()


def render_pptx(lines: Iterable[str], trip: Trip | None = None) -> bytes:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    all_lines = list(lines)
    if trip and trip.days:
        cover = presentation.slides.add_slide(presentation.slide_layouts[6])
        title = cover.shapes.add_textbox(Inches(0.65), Inches(0.3), Inches(12), Inches(0.65))
        title.text_frame.text = trip.title
        title.text_frame.paragraphs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].font.bold = True
        cover.shapes.add_picture(_render_route_map(trip, 1600, 720), Inches(0.65), Inches(1.15), width=Inches(12), height=Inches(5.4))
        highlights = presentation.slides.add_slide(presentation.slide_layouts[6])
        heading = highlights.shapes.add_textbox(Inches(0.65), Inches(0.3), Inches(12), Inches(0.6))
        heading.text_frame.text = "景点、餐饮与住宿详情"
        heading.text_frame.paragraphs[0].font.size = Pt(27)
        heading.text_frame.paragraphs[0].font.bold = True
        highlights.shapes.add_picture(
            _render_activity_board(trip, 1600, 720),
            Inches(0.65), Inches(1.1), width=Inches(12), height=Inches(5.4),
        )
        for day in trip.days:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            title = slide.shapes.add_textbox(Inches(0.65), Inches(0.3), Inches(12), Inches(0.55))
            title.text_frame.text = f"第 {day.day_index} 天  ·  {day.date}  {day.title}"
            title.text_frame.paragraphs[0].font.size = Pt(25)
            title.text_frame.paragraphs[0].font.bold = True
            subtitle = slide.shapes.add_textbox(Inches(0.68), Inches(0.85), Inches(12), Inches(0.35))
            subtitle.text_frame.text = f"{day.total_distance_km:g} km  ·  驾驶 {day.total_drive_minutes // 60}h{day.total_drive_minutes % 60}m  ·  {len(day.activities)} 项安排"
            subtitle.text_frame.paragraphs[0].font.size = Pt(12)
            slide.shapes.add_picture(
                _render_route_map(trip, 1600, 300, days=[day]),
                Inches(0.65), Inches(1.2), width=Inches(12), height=Inches(1.8),
            )
            slide.shapes.add_picture(
                _render_activity_board(trip, 1600, 620, activities=day.activities),
                Inches(0.65), Inches(3.1), width=Inches(12), height=Inches(2.85),
            )
            body = slide.shapes.add_textbox(Inches(0.85), Inches(6.1), Inches(11.5), Inches(0.9))
            body.text_frame.text = "\n".join(
                f"{stage.planned_start:%H:%M}–{stage.planned_end:%H:%M}  {stage.origin.name} → {stage.destination.name}  ·  {stage.mode}  ·  {stage.distance_km:g} km"
                for stage in day.stages[:4]
            )
            for paragraph in body.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
        # Visual slides already contain the structured report; avoid appending
        # the old wall of raw Markdown lines after them.
        output = BytesIO()
        presentation.save(output)
        return output.getvalue()
    chunks = [all_lines[index:index + 28] for index in range(0, len(all_lines), 28)] or [[]]
    for slide_index, chunk in enumerate(chunks):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.55))
        title.text_frame.text = chunk[0] if slide_index == 0 and chunk else f"RoadMan 行程 · 第 {slide_index + 1} 页"
        title.text_frame.paragraphs[0].font.size = Pt(26)
        title.text_frame.paragraphs[0].font.bold = True
        body = slide.shapes.add_textbox(Inches(0.75), Inches(1.1), Inches(11.8), Inches(5.9))
        body.text_frame.text = "\n".join(chunk[1:] if slide_index == 0 else chunk)
        for paragraph in body.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def render_long_image(lines: Iterable[str], trip: Trip | None = None) -> bytes:
    all_lines = list(lines)
    width, line_height, padding = 1600, 34, 56
    visual_height = 1080 if trip and trip.days else 0
    if trip and trip.days:
        visual_height += len(trip.days) * 900
    height = max(900, padding * 2 + visual_height + line_height * max(1, len(all_lines)))
    image = Image.new("RGB", (width, height), "#f4f8ff")
    draw = ImageDraw.Draw(image)
    font = _load_font(24)
    title_font = _load_font(38)
    y = padding
    if trip and trip.days:
        draw.rounded_rectangle((padding, y, width - padding, y + 66), radius=22, fill="#10213e")
        draw.text((padding + 28, y + 15), trip.title[:42], fill="#ffffff", font=title_font)
        y += 92
        route = Image.open(_render_route_map(trip, width - padding * 2, 560)).convert("RGB")
        image.paste(route, (padding, y))
        y += 590
        activity_board = Image.open(
            _render_activity_board(trip, width - padding * 2, 390)
        ).convert("RGB")
        image.paste(activity_board, (padding, y))
        y += 420
        for day in trip.days:
            draw.text((padding, y), f"第 {day.day_index} 天 · {day.date} · {day.title}", fill="#10213e", font=_load_font(30))
            y += 48
            day_route = Image.open(
                _render_route_map(trip, width - padding * 2, 280, days=[day])
            ).convert("RGB")
            image.paste(day_route, (padding, y))
            y += 300
            day_board = Image.open(
                _render_activity_board(trip, width - padding * 2, 400, activities=day.activities)
            ).convert("RGB")
            image.paste(day_board, (padding, y))
            y += 430
            draw.rounded_rectangle((padding, y, width - padding, y + 12 + max(34, len(day.stages) * 28)), radius=16, fill="#ffffff", outline="#dce7f4", width=2)
            text_y = y + 16
            for stage in day.stages:
                text = f"{stage.planned_start:%H:%M}–{stage.planned_end:%H:%M}   {stage.origin.name} → {stage.destination.name}   {stage.mode} · {stage.distance_km:g} km"
                draw.text((padding + 20, text_y), text[:120], fill="#385171", font=font)
                text_y += 28
            y = text_y + 24
    else:
        for index, line in enumerate(all_lines):
            draw.text((padding, y), line[:105], fill="#10213e", font=title_font if index == 0 else font)
            y += 48 if index == 0 else line_height
    if trip and trip.days:
        image = image.crop((0, 0, width, min(height, y + padding)))
    output = BytesIO()
    image.save(output, format="PNG", optimize=True)
    return output.getvalue()


def _render_route_map(
    trip: Trip,
    width: int,
    height: int,
    *,
    days: list | None = None,
) -> BytesIO:
    """Render a route overview with an AMap geographic base when available."""
    selected_days = days if days is not None else trip.days
    image = Image.new("RGB", (width, height), "#eef5fc")
    draw = ImageDraw.Draw(image)
    margin = max(34, width // 28)
    for x in range(margin, width - margin, max(80, width // 12)):
        draw.line((x, margin, x, height - margin), fill="#dce8f3", width=1)
    for y in range(margin, height - margin, max(70, height // 8)):
        draw.line((margin, y, width - margin, y), fill="#dce8f3", width=1)

    routes: list[tuple[str, list[tuple[float, float]]]] = []
    for day in selected_days:
        for stage in day.stages:
            points = [
                (point.longitude, point.latitude)
                for segment in stage.route_segments
                for point in segment.coordinates
            ]
            if len(points) >= 2:
                routes.append((stage.mode, points))
    all_points = [point for _, points in routes for point in points]
    if not all_points:
        output = BytesIO()
        image.save(output, format="PNG")
        output.seek(0)
        return output
    min_lon, max_lon = min(p[0] for p in all_points), max(p[0] for p in all_points)
    min_lat, max_lat = min(p[1] for p in all_points), max(p[1] for p in all_points)
    lon_span = max(max_lon - min_lon, 0.01)
    lat_span = max(max_lat - min_lat, 0.01)

    static_map = _fetch_amap_static_map(
        routes,
        width,
        height,
        center=((min_lon + max_lon) / 2, (min_lat + max_lat) / 2),
        span=max(lon_span, lat_span),
    )
    if static_map is not None:
        # AMap already paints the true route geometry over satellite tiles.
        # Keep the local fallback renderer below for offline/export tests.
        output = BytesIO()
        static_map.save(output, format="PNG", optimize=True)
        output.seek(0)
        return output

    def project(point: tuple[float, float]) -> tuple[int, int]:
        x = margin + int((point[0] - min_lon) / lon_span * (width - margin * 2))
        y = height - margin - int((point[1] - min_lat) / lat_span * (height - margin * 2))
        return x, y

    colors = {"driving": "#1777e8", "transit": "#20a879", "walking": "#f2a51a", "riding": "#e6a11d"}
    for mode, points in routes:
        projected = [project(point) for point in points]
        draw.line(projected, fill="#ffffff", width=max(10, width // 105), joint="curve")
        draw.line(projected, fill=colors.get(mode, "#72849a"), width=max(5, width // 190), joint="curve")
    numbered: list[tuple[int, int]] = []
    for _, points in routes:
        for point in (points[0], points[-1]):
            pixel = project(point)
            if not numbered or min((pixel[0]-x)**2 + (pixel[1]-y)**2 for x, y in numbered) > 500:
                numbered.append(pixel)
    font = _load_font(max(16, width // 65))
    for index, (x, y) in enumerate(numbered, 1):
        radius = max(13, width // 85)
        draw.ellipse((x-radius, y-radius, x+radius, y+radius), fill="#ffffff", outline="#176fdf", width=3)
        text = str(index)
        box = draw.textbbox((0, 0), text, font=font)
        draw.text((x-(box[2]-box[0])/2, y-(box[3]-box[1])/2-1), text, fill="#145fbf", font=font)
    caption_font = _load_font(max(12, width // 92))
    draw.rounded_rectangle((margin, height - 34, min(width - margin, margin + width // 2), height - 8), radius=10, fill="#ffffff")
    draw.text(
        (margin + 10, height - 30),
        "真实道路点列 · 路线来源：高德 WebService · 景点来源详见活动卡",
        fill="#48617f", font=caption_font,
    )
    output = BytesIO()
    image.save(output, format="PNG", optimize=True)
    output.seek(0)
    return output


def _fetch_amap_static_map(
    routes: list[tuple[str, list[tuple[float, float]]]],
    width: int,
    height: int,
    *,
    center: tuple[float, float],
    span: float,
) -> Image.Image | None:
    """Fetch an AMap geographic map with route paths, degrading to local map.

    Static-map rendering is best effort: an export must remain usable when
    the map provider is unavailable, while configured AMap credentials should
    produce a geographic background instead of a bare grid.
    """
    key = get_settings().amap_webservice_key
    if not key:
        return None
    zoom = (
        5 if span > 8 else 6 if span > 4 else 7 if span > 2 else
        8 if span > 1 else 9 if span > 0.5 else 10 if span > 0.2 else
        11 if span > 0.1 else 13
    )
    path_specs: list[str] = []
    colors = {
        "driving": "0x1777e8",
        "transit": "0x20a879",
        "walking": "0xf2a51a",
        "riding": "0xe6a11d",
    }
    for mode, points in routes:
        sampled = _sample_route_points(points, 180)
        if len(sampled) < 2:
            continue
        encoded_points = ";".join(f"{lon:.6f},{lat:.6f}" for lon, lat in sampled)
        # AMap's web-service syntax is positional:
        # weight,color,transparency,fillcolor,fillTransparency:points.
        path_specs.append(f"6,{colors.get(mode, '0x72849a')},1,,:{encoded_points}")
    if not path_specs:
        return None
    # The static-map endpoint accepts at most four paths.  Stages in a trip
    # are contiguous, so when there are more segments keep one continuous
    # geometry rather than silently dropping later stages.
    if len(path_specs) > 4:
        merged_points: list[tuple[float, float]] = []
        for _, points in routes:
            sampled = _sample_route_points(points, 180)
            if merged_points and sampled and merged_points[-1] == sampled[0]:
                sampled = sampled[1:]
            merged_points.extend(sampled)
        if len(merged_points) >= 2:
            encoded_points = ";".join(f"{lon:.6f},{lat:.6f}" for lon, lat in _sample_route_points(merged_points, 500))
            path_specs = [f"6,0x1777e8,1,,:{encoded_points}"]
    map_width = min(1024, max(420, int(width)))
    map_height = min(1024, max(280, int(height)))
    try:
        response = httpx.get(
            "https://restapi.amap.com/v3/staticmap",
            params={
                "location": f"{center[0]:.6f},{center[1]:.6f}",
                "zoom": zoom,
                "size": f"{map_width}*{map_height}",
                "scale": 1,
                "paths": "|".join(path_specs),
                "key": key,
            },
            timeout=2.5,
            follow_redirects=True,
        )
        response.raise_for_status()
        with Image.open(BytesIO(response.content)) as loaded:
            return ImageOps.fit(loaded.convert("RGB"), (width, height), method=Image.Resampling.LANCZOS)
    except (httpx.HTTPError, OSError, ValueError):
        return None


def _sample_route_points(points: list[tuple[float, float]], limit: int) -> list[tuple[float, float]]:
    if len(points) <= limit:
        return points
    stride = (len(points) - 1) / (limit - 1)
    return [points[round(index * stride)] for index in range(limit)]


def _render_activity_board(
    trip: Trip,
    width: int,
    height: int,
    *,
    activities: list | None = None,
) -> BytesIO:
    """Render the curated POI/hotel/meal cards used by every rich export."""
    image = Image.new("RGB", (width, height), "#f7faff")
    draw = ImageDraw.Draw(image)
    title_font = _load_font(max(20, width // 50))
    body_font = _load_font(max(14, width // 78))
    small_font = _load_font(max(12, width // 95))
    seen_image_urls: set[str] = set()
    draw.text((28, 20), "智能体精选行程安排", fill="#10213e", font=title_font)
    activities = sorted(
        activities
        if activities is not None
        else [
            activity
            for day in trip.days
            for activity in day.activities
            if activity.type in {"attraction", "hotel", "meal"}
        ],
        key=lambda item: ({"attraction": 0, "hotel": 1, "meal": 2}.get(item.type, 3), item.planned_start),
    )[:4]
    if not activities:
        draw.text((28, 78), "暂无可展示的景点、餐饮或住宿安排", fill="#66758f", font=body_font)
    else:
        gap = 24
        card_width = (width - 56 - gap * (len(activities) - 1)) // len(activities)
        card_top, card_bottom = 72, height - 24
        for index, activity in enumerate(activities):
            left = 28 + index * (card_width + gap)
            right = left + card_width
            draw.rounded_rectangle(
                (left, card_top, right, card_bottom), radius=18,
                fill="#ffffff", outline="#dce7f4", width=2,
            )
            photo = _activity_image(activity, seen_image_urls)
            photo_height = max(80, int((card_bottom - card_top) * 0.43))
            if photo:
                fitted = ImageOps.fit(photo.convert("RGB"), (card_width - 4, photo_height))
                image.paste(fitted, (left + 2, card_top + 2))
            text_y = card_top + photo_height + 14
            draw.text((left + 14, text_y), activity.place.name[:16], fill="#10213e", font=body_font)
            text_y += max(24, width // 60)
            time_text = f"{activity.planned_start:%m月%d日 %H:%M}–{activity.planned_end:%H:%M}"
            draw.text((left + 14, text_y), time_text, fill="#3e5c82", font=small_font)
            text_y += max(20, width // 78)
            check_text = {
                "required": "需预约",
                "recommended": "建议预约",
                "not_required": "无需预约",
                "unknown": "预约待核查",
            }.get(activity.reservation_status, "预约待核查")
            if activity.risk_tags:
                check_text += " · " + "、".join(activity.risk_tags[:2])
            draw.text((left + 14, text_y), check_text[:28], fill="#9a6a24", font=small_font)
            text_y += max(18, width // 88)
            note = activity.description or activity.user_note or next(
                (source.title for source in activity.source_records if source.title),
                "详情见来源链接",
            )
            for line in _wrap(note, max(8, card_width // max(12, width // 85)))[:2]:
                draw.text((left + 14, text_y), line, fill="#66758f", font=small_font)
                text_y += max(18, width // 88)
    output = BytesIO()
    image.save(output, format="PNG", optimize=True)
    output.seek(0)
    return output


def _download_activity_image(url: str | None) -> Image.Image | None:
    if not url:
        return None
    parsed = urlparse(url)
    if parsed.scheme not in {"http", "https"} or not parsed.hostname:
        return None
    try:
        response = httpx.get(url, timeout=2.5, follow_redirects=True)
        response.raise_for_status()
        if len(response.content) > 5 * 1024 * 1024:
            return None
        return Image.open(BytesIO(response.content)).convert("RGB")
    except (httpx.HTTPError, OSError, ValueError):
        return None


def _activity_image(activity, seen_image_urls: set[str]) -> Image.Image:
    """Return a real, non-repeated image or a distinct visual placeholder."""
    url = str(activity.image_url or "").strip()
    if url and url not in seen_image_urls:
        image = _download_activity_image(url)
        if image is not None:
            seen_image_urls.add(url)
            return image
    return _make_activity_placeholder(activity)


def _make_activity_placeholder(activity, width: int = 960, height: int = 600) -> Image.Image:
    """Create a stable card visual so missing/reused provider images differ."""
    seed = sum((index + 1) * ord(char) for index, char in enumerate(activity.place.name))
    palettes = ((35, 119, 232), (20, 155, 132), (116, 87, 232), (225, 148, 37), (34, 104, 145))
    first = palettes[seed % len(palettes)]
    second = palettes[(seed // 7 + 2) % len(palettes)]
    image = Image.new("RGB", (width, height), first)
    draw = ImageDraw.Draw(image)
    for y in range(height):
        ratio = y / max(1, height - 1)
        color = tuple(int(first[index] * (1 - ratio) + second[index] * ratio) for index in range(3))
        draw.line((0, y, width, y), fill=color)
    draw.ellipse((width * 0.62, -height * 0.25, width * 1.15, height * 0.55), fill=(245, 252, 255))
    draw.ellipse((-width * 0.2, height * 0.55, width * 0.45, height * 1.2), fill=(225, 241, 255))
    label = _activity_label(activity.type)
    title = activity.place.name[:14]
    label_font = _load_font(max(26, width // 25))
    title_font = _load_font(max(22, width // 34))
    draw.rounded_rectangle((38, 38, 38 + max(120, len(label) * 42), 96), radius=28, fill="#ffffff")
    draw.text((64, 49), label, fill="#173b68", font=label_font)
    draw.text((48, height - 130), title, fill="#ffffff", font=title_font)
    draw.text((48, height - 82), "RoadMan · 来源待确认", fill="#eaf4ff", font=_load_font(max(16, width // 58)))
    return image


def _wrap(value: str, width: int) -> list[str]:
    return [value[index:index + width] for index in range(0, len(value), width)] or [""]


def _load_font(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    for candidate in (
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "C:/Windows/Fonts/msyh.ttc",
    ):
        if Path(candidate).is_file():
            try:
                return ImageFont.truetype(candidate, size)
            except OSError:
                continue
    return ImageFont.load_default()
